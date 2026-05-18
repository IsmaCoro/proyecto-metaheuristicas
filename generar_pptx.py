"""
Script para generar la presentacion en PowerPoint (.pptx)
Ejecutar:  pip install python-pptx  &&  python generar_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colores del tema
FONDO = RGBColor(0x0a, 0x0a, 0x1a)
ACENTO = RGBColor(0x66, 0x7e, 0xea)
BLANCO = RGBColor(0xff, 0xff, 0xff)
GRIS = RGBColor(0xaa, 0xaa, 0xbb)
ROJO = RGBColor(0xE0, 0x6C, 0x75)
CYAN = RGBColor(0x56, 0xB6, 0xC2)
VERDE = RGBColor(0x56, 0xD6, 0x8C)
MORADO = RGBColor(0x76, 0x4B, 0xA2)
AMARILLO = RGBColor(0xD1, 0x9A, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text(slide, left, top, width, height, text, size=18, bold=False, color=BLANCO, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_para(tf, text, size=18, bold=False, color=GRIS, align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return p

# ═══════════════════════════════════════════
# SLIDE 1: Portada
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, FONDO)
add_text(s, 1, 0.8, 11, 0.6, "PROYECTO FINAL — ALGORITMOS METAHEURÍSTICOS", 16, True, ACENTO, PP_ALIGN.CENTER)
add_text(s, 1, 1.8, 11, 1.2, "Optimización Inteligente de\nRecursos en la Nube", 40, True, BLANCO, PP_ALIGN.CENTER)
add_text(s, 1, 3.5, 11, 0.5, "Jonathan Ismael Corona Mendez", 20, False, GRIS, PP_ALIGN.CENTER)
add_text(s, 1, 4.1, 11, 0.5, "Universidad de Guadalajara — CUCosta, Puerto Vallarta", 16, False, GRIS, PP_ALIGN.CENTER)
add_text(s, 1, 4.6, 11, 0.5, "Ingeniería en Computación, 8.° Semestre — 18 de mayo de 2026", 14, False, GRIS, PP_ALIGN.CENTER)
tf = add_text(s, 2.5, 5.5, 8, 0.5, "Algoritmo Genético  ·  Colonia de Hormigas  ·  Híbrido GA-ACO  ·  Wilcoxon", 13, False, ACENTO, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 2: El Problema
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, FONDO)
add_text(s, 0.8, 0.4, 5, 0.4, "PLANTEAMIENTO", 14, True, ACENTO)
add_text(s, 0.8, 0.9, 11, 0.7, "¿Cuál es el problema?", 32, True, BLANCO)

add_text(s, 0.8, 1.9, 5.5, 0.5, "Asignación de Recursos en la Nube", 20, True, BLANCO)
tf = add_text(s, 0.8, 2.5, 5.5, 3, "Un datacenter con 6 servidores heterogéneos recibe tareas con distintos requerimientos de CPU y RAM.", 16, False, GRIS)
add_para(tf, "")
add_para(tf, "▸ Minimizar desbalance de carga", 15, False, GRIS)
add_para(tf, "▸ Reducir consumo energético", 15, False, GRIS)
add_para(tf, "▸ Evitar cuellos de botella", 15, False, GRIS)
add_para(tf, "▸ Respetar restricciones de capacidad", 15, False, GRIS)

add_text(s, 7, 1.9, 5.5, 0.5, "¿Por qué es difícil?", 20, True, BLANCO)
tf = add_text(s, 7, 2.5, 5.5, 1, "Es un problema NP-duro. Con 20 tareas y 6 servidores:", 16, False, GRIS)
add_para(tf, "")
add_para(tf, "6²⁰ = 3.6 × 10¹⁵ combinaciones", 28, True, ACENTO, PP_ALIGN.CENTER)
add_para(tf, "")
add_para(tf, "No se pueden probar todas.\nSe necesitan metaheurísticas.", 16, False, GRIS)

# ═══════════════════════════════════════════
# SLIDE 3: Fitness
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, RGBColor(0x0d, 0x0d, 0x2b))
add_text(s, 0.8, 0.4, 5, 0.4, "FUNCIÓN OBJETIVO", 14, True, ACENTO)
add_text(s, 0.8, 0.9, 11, 0.7, "Fitness — Menor = Mejor", 32, True, BLANCO)

add_text(s, 1.5, 2.0, 10, 0.6, "f(x) = 0.4·Balance + 0.3·Energía/100 + 0.3·Tiempo·10 + Penalización", 18, True, RGBColor(0xc0, 0xc8, 0xff), PP_ALIGN.CENTER)

add_text(s, 0.8, 3.2, 2.8, 0.4, "⚖️ Balance (40%)", 16, True, BLANCO)
add_text(s, 0.8, 3.7, 2.8, 0.8, "Desviación estándar del uso de CPU entre servidores", 13, False, GRIS)

add_text(s, 3.8, 3.2, 2.8, 0.4, "⚡ Energía (30%)", 16, True, BLANCO)
add_text(s, 3.8, 3.7, 2.8, 0.8, "Watts totales del datacenter. Menos = mejor", 13, False, GRIS)

add_text(s, 6.8, 3.2, 2.8, 0.4, "⏱️ Tiempo (30%)", 16, True, BLANCO)
add_text(s, 6.8, 3.7, 2.8, 0.8, "Servidor más saturado = cuello de botella", 13, False, GRIS)

add_text(s, 9.8, 3.2, 2.8, 0.4, "🚫 Penalización", 16, True, BLANCO)
add_text(s, 9.8, 3.7, 2.8, 0.8, "Castigo por exceder CPU o RAM del servidor", 13, False, GRIS)

add_text(s, 0.8, 5.2, 11, 0.5, "Ejemplo:  Distribución manual → Fitness = 85   |   Optimizada → Fitness = 8", 16, False, VERDE, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 4: Algoritmo Genético
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, FONDO)
add_text(s, 0.8, 0.4, 5, 0.4, "METAHEURÍSTICA 1", 14, True, ACENTO)
add_text(s, 0.8, 0.9, 11, 0.7, "Algoritmo Genético (GA)", 32, True, BLANCO)

add_text(s, 0.8, 1.9, 5.5, 0.5, "Inspiración: Evolución de Darwin", 18, True, BLANCO)
tf = add_text(s, 0.8, 2.5, 5.5, 3.5, "Una población de soluciones evoluciona generación tras generación.", 15, False, GRIS)
add_para(tf, "")
add_para(tf, "▸ Selección por torneo (k=5)", 15, False, GRIS)
add_para(tf, "▸ Cruce uniforme (p=0.85)", 15, False, GRIS)
add_para(tf, "▸ Mutación aleatoria (p=0.15)", 15, False, GRIS)
add_para(tf, "▸ Elitismo: top 3 pasan directo", 15, False, GRIS)
add_para(tf, "▸ 60 individuos × 150 generaciones", 15, False, GRIS)

add_text(s, 7, 1.9, 5.5, 0.5, "Representación", 18, True, BLANCO)
add_text(s, 7, 2.5, 5.5, 0.8, "Individuo = [2, 0, 5, 1, 3, 0, 4, 2, ...]\nGen[i] = servidor asignado a tarea i", 14, False, VERDE)
add_text(s, 7, 3.8, 5.5, 0.4, "Características", 18, True, BLANCO)
tf = add_text(s, 7, 4.3, 5.5, 1.5, "✓ Buena exploración global", 15, False, VERDE)
add_para(tf, "✗ Se estanca en óptimos locales", 15, False, ROJO)

# ═══════════════════════════════════════════
# SLIDE 5: ACO
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, FONDO)
add_text(s, 0.8, 0.4, 5, 0.4, "METAHEURÍSTICA 2", 14, True, ACENTO)
add_text(s, 0.8, 0.9, 11, 0.7, "Colonia de Hormigas (ACO)", 32, True, BLANCO)

add_text(s, 0.8, 1.9, 5.5, 0.5, "Inspiración: Feromonas de hormigas", 18, True, BLANCO)
tf = add_text(s, 0.8, 2.5, 5.5, 1.5, "Las hormigas depositan feromonas en los mejores caminos. Las siguientes hormigas siguen esas señales.", 15, False, GRIS)
add_para(tf, "")
add_para(tf, "P(i→j) = (τ^α · η^β) / Σ(τ^α · η^β)", 16, True, RGBColor(0xc0, 0xc8, 0xff))
add_para(tf, "τ=feromona  η=capacidad libre  α=1.0  β=2.0", 12, False, GRIS)

add_text(s, 7, 1.9, 5.5, 0.5, "Parámetros", 18, True, BLANCO)
tf = add_text(s, 7, 2.5, 5.5, 2, "▸ 30 hormigas × 150 iteraciones", 15, False, GRIS)
add_para(tf, "▸ Evaporación: ρ = 0.3 (30%)", 15, False, GRIS)
add_para(tf, "▸ Depósito: Q/fitness", 15, False, GRIS)
add_para(tf, "")
add_para(tf, "Características", 18, True, BLANCO)
add_para(tf, "✓ Buena intensificación", 15, False, VERDE)
add_para(tf, "✗ Convergencia lenta al inicio", 15, False, CYAN)

# ═══════════════════════════════════════════
# SLIDE 6: Híbrido
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, RGBColor(0x0d, 0x0d, 0x2b))
add_text(s, 0.8, 0.4, 8, 0.4, "METAHEURÍSTICA 3 — APORTACIÓN ORIGINAL", 14, True, ACENTO)
add_text(s, 0.8, 0.9, 11, 0.7, "Híbrido GA-ACO", 32, True, BLANCO)

add_text(s, 0.8, 1.9, 5.5, 0.5, "Combinación de ambos paradigmas", 18, True, BLANCO)
tf = add_text(s, 0.8, 2.5, 5.5, 3.5, "▸ 70% GA: cruce + mutación guiada por feromonas", 15, False, GRIS)
add_para(tf, "▸ 30% ACO: construcción probabilística", 15, False, GRIS)
add_para(tf, "▸ Mutación NO aleatoria → usa feromonas", 15, False, GRIS)
add_para(tf, "▸ Búsqueda local cada 20 generaciones", 15, False, GRIS)
add_para(tf, "▸ Refinamiento final del mejor global", 15, False, GRIS)
add_para(tf, "")
add_para(tf, "Exploración (GA) + Intensificación (ACO)", 16, True, ACENTO, PP_ALIGN.CENTER)
add_para(tf, "+ Refinamiento (Búsqueda Local)", 16, True, ACENTO, PP_ALIGN.CENTER)

add_text(s, 7, 1.9, 5.5, 0.5, "Pseudocódigo", 18, True, BLANCO)
code = """P ← 1 greedy + 79 aleatorias
τ ← feromonas uniformes

PARA gen = 1..150:
  evaluar(P)
  actualizar feromonas(élite)
  nueva ← top 5 (elitismo)
  PARA 70% restante:    // GA
    hijo ← cruce(torneo, torneo)
    mutar_con_feromonas(hijo, τ)
  PARA 30% restante:    // ACO
    sol ← construir(τ, heurística)
  SI gen % 20 = 0:
    busqueda_local(mejor)

busqueda_local(mejor_global)
RETORNAR mejor_global"""
add_text(s, 7, 2.5, 5.5, 4.5, code, 11, False, VERDE)

# ═══════════════════════════════════════════
# SLIDE 7: Resultados
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, FONDO)
add_text(s, 0.8, 0.4, 5, 0.4, "EXPERIMENTACIÓN", 14, True, ACENTO)
add_text(s, 0.8, 0.9, 11, 0.7, "Resultados de Optimización", 32, True, BLANCO)

# Fitness boxes
for i, (lbl, val, col) in enumerate([("Greedy", "~25", AMARILLO), ("GA", "~22", ROJO), ("ACO", "~18", CYAN), ("Híbrido", "~12", ACENTO)]):
    x = 1.5 + i * 2.7
    add_text(s, x, 2.2, 2.2, 0.8, val, 36, True, col, PP_ALIGN.CENTER)
    add_text(s, x, 3.0, 2.2, 0.4, lbl, 14, False, GRIS, PP_ALIGN.CENTER)

add_text(s, 1, 4.0, 5, 0.4, "Antes (usuario)", 18, True, ROJO, PP_ALIGN.CENTER)
add_text(s, 1, 4.5, 5, 1.5, "S0: 120% ← SATURADO\nS1: 12%\nS2: 0%\nS3: 88%", 15, False, GRIS, PP_ALIGN.CENTER)

add_text(s, 7, 4.0, 5, 0.4, "Después (Híbrido)", 18, True, VERDE, PP_ALIGN.CENTER)
add_text(s, 7, 4.5, 5, 1.5, "S0: 25%\nS1: 30%\nS2: 28%\nS3: 17%", 15, False, VERDE, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 8: Wilcoxon
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, FONDO)
add_text(s, 0.8, 0.4, 5, 0.4, "VALIDACIÓN ESTADÍSTICA", 14, True, ACENTO)
add_text(s, 0.8, 0.9, 11, 0.7, "Prueba de Wilcoxon — 10 Corridas", 32, True, BLANCO)

# Stats table header
y = 2.0
for i, h in enumerate(["Algoritmo", "Media", "Mejor", "Peor", "Desv.Std"]):
    add_text(s, 1 + i*2.2, y, 2, 0.4, h, 14, True, ACENTO, PP_ALIGN.CENTER)
for j, (alg, col, vals) in enumerate([("GA", ROJO, ["~24.4","~21","~28","~2.5"]), ("ACO", CYAN, ["~19.9","~18","~22","~1.3"]), ("Híbrido", ACENTO, ["~12.7","~11","~15","~1.4"])]):
    yy = y + 0.5 + j*0.45
    add_text(s, 1, yy, 2, 0.4, alg, 14, True, col, PP_ALIGN.CENTER)
    for k, v in enumerate(vals):
        add_text(s, 3.2 + k*2.2, yy, 2, 0.4, v, 14, False, GRIS, PP_ALIGN.CENTER)

# Wilcoxon table
y2 = 4.2
for i, h in enumerate(["Par", "p-valor", "Significativo", "Mejor"]):
    add_text(s, 1.5 + i*2.5, y2, 2.3, 0.4, h, 14, True, ACENTO, PP_ALIGN.CENTER)
for j, (par, mejor) in enumerate([("GA vs ACO", "ACO"), ("GA vs Híbrido", "Híbrido"), ("ACO vs Híbrido", "Híbrido")]):
    yy = y2 + 0.5 + j*0.45
    add_text(s, 1.5, yy, 2.3, 0.4, par, 13, False, GRIS, PP_ALIGN.CENTER)
    add_text(s, 4, yy, 2.3, 0.4, "<0.05", 13, False, GRIS, PP_ALIGN.CENTER)
    add_text(s, 6.5, yy, 2.3, 0.4, "Sí ✓", 13, False, VERDE, PP_ALIGN.CENTER)
    add_text(s, 9, yy, 2.3, 0.4, mejor, 13, True, ACENTO, PP_ALIGN.CENTER)

add_text(s, 1, 6.2, 11, 0.5, "Con p < 0.05, las diferencias NO son producto del azar.", 16, False, BLANCO, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 9: Sensibilidad
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, FONDO)
add_text(s, 0.8, 0.4, 5, 0.4, "ANÁLISIS DE SENSIBILIDAD", 14, True, ACENTO)
add_text(s, 0.8, 0.9, 11, 0.7, "¿Son los mejores parámetros?", 32, True, BLANCO)

add_text(s, 0.8, 1.9, 5.5, 0.5, "Tamaño de Población", 20, True, BLANCO)
tf = add_text(s, 0.8, 2.5, 5.5, 2.5, "Se probó con 20, 40, 60, 80, 100 individuos (iteraciones fijas = 80)", 14, False, GRIS)
add_para(tf, "")
add_para(tf, "▸ 20 → fitness alto (poca diversidad)", 15, False, GRIS)
add_para(tf, "▸ 60 → punto de equilibrio ✓", 15, True, VERDE)
add_para(tf, "▸ 80-100 → rendimiento decreciente", 15, False, GRIS)

add_text(s, 7, 1.9, 5.5, 0.5, "Número de Iteraciones", 20, True, BLANCO)
tf = add_text(s, 7, 2.5, 5.5, 2.5, "Se probó con 30, 60, 100, 150, 200 (población fija = 60)", 14, False, GRIS)
add_para(tf, "")
add_para(tf, "▸ 30 → no converge", 15, False, GRIS)
add_para(tf, "▸ 100 → se estabiliza", 15, False, GRIS)
add_para(tf, "▸ 150 → valor óptimo ✓", 15, True, VERDE)
add_para(tf, "▸ 200 → mejora insignificante", 15, False, GRIS)

add_text(s, 1.5, 5.5, 10, 0.6, "Parámetros óptimos:  60 individuos  ×  150 iteraciones", 22, True, ACENTO, PP_ALIGN.CENTER)
add_text(s, 1.5, 6.2, 10, 0.4, "Equilibrio entre calidad de solución y costo computacional", 14, False, GRIS, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════
# SLIDE 10: Conclusiones
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, RGBColor(0x0d, 0x0d, 0x2b))
add_text(s, 0.8, 0.4, 5, 0.4, "CIERRE", 14, True, ACENTO)
add_text(s, 0.8, 0.9, 11, 0.7, "Conclusiones", 32, True, BLANCO)

conclusions = [
    ("1. Superan al baseline", "Las 3 metaheurísticas superan tanto la asignación manual como el greedy determinista."),
    ("2. El Híbrido gana", "Combina exploración (GA) + intensificación (ACO) + refinamiento local = mejor fitness."),
    ("3. Validación rigurosa", "Wilcoxon con p<0.05 confirma que las diferencias son estadísticamente significativas."),
    ("4. Parámetros calibrados", "Análisis de sensibilidad identifica 60 individuos y 150 iteraciones como punto óptimo."),
]
for i, (titulo, desc) in enumerate(conclusions):
    x = 0.8 + (i % 2) * 6
    y = 2.0 + (i // 2) * 2.2
    add_text(s, x, y, 5.5, 0.5, titulo, 18, True, BLANCO)
    add_text(s, x, y + 0.5, 5.5, 1, desc, 14, False, GRIS)

add_text(s, 0.8, 6.0, 5, 0.4, "Trabajo Futuro", 16, True, GRIS)
add_text(s, 0.8, 6.4, 11, 0.4, "NSGA-II (Pareto)  ·  Migración dinámica  ·  Escalabilidad  ·  ML predictivo  ·  Datos reales", 13, False, GRIS)

# ═══════════════════════════════════════════
# SLIDE 11: Demo
# ═══════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, FONDO)
add_text(s, 1, 1.5, 11, 0.5, "DEMOSTRACIÓN EN VIVO", 16, True, ACENTO, PP_ALIGN.CENTER)
add_text(s, 1, 2.5, 11, 1, "Demo Time", 56, True, BLANCO, PP_ALIGN.CENTER)
add_text(s, 1, 4.0, 11, 0.5, "Abran su celular y entren a la URL", 20, False, GRIS, PP_ALIGN.CENTER)
add_text(s, 2.5, 5.0, 8, 0.8, "http://TU_IP:5000", 36, True, ACENTO, PP_ALIGN.CENTER)
add_text(s, 2.5, 5.9, 8, 0.5, "Agreguen tareas a los servidores", 14, False, GRIS, PP_ALIGN.CENTER)

# Guardar
prs.save("presentacion.pptx")
print("presentacion.pptx generado correctamente!")
